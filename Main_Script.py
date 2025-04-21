# Importing Libraries

import os
import faiss
import numpy as np
import pdfplumber
import pandas as pd
import json
import easyocr
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from groq import Groq


# Setting up API key and client

MY_API_KEY = "enter_your_API_key_here"
client = Groq(api_key = MY_API_KEY)


# Setting up embedding model 

embedding_model = SentenceTransformer('all-MiniLM-L6-v2')


# FAISS Setup 
EMBED_DIM = 384
index = faiss.IndexFlatL2(EMBED_DIM)
document_chunks = []


# Setting up document loaders 

def load_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Page {i+1} ---\n{page_text}"
    return text

def load_docx(file_path):
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for i, slide in enumerate(prs.slides):
        lines = [shape.text.strip() for shape in slide.shapes
                 if hasattr(shape, "text") and shape.text.strip()]
        if lines:
            out += f"\n--- Slide {i+1} ---\n" + "\n".join(lines)
    return text

def load_excel(file_path):
    return pd.read_excel(file_path).to_string(index=False)

def load_csv(file_path):
    return pd.read_csv(file_path).to_string(index=False)

def load_json(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return json.dumps(json.load(f), indent=2)

def load_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def load_image(file_path):
    reader = easyocr.Reader(['en'])
    result = reader.readtext(file_path)
    text = "\n".join([detection[1] for detection in result])
    return text


# Chunking & Embedding Setup

def split_text(text, chunk_size=300):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def get_embedding(text):
    return embedding_model.encode(text, convert_to_numpy=True)

def add_documents(chunks):
    for chunk in chunks:
        vec = get_embedding(chunk)
        index.add(np.expand_dims(vec, axis=0))
        document_chunks.append(chunk)


# Retrieval

def retrieve_context(query, k=3):
    qv = get_embedding(query)
    _, idxs = index.search(np.expand_dims(qv, axis=0), k)
    return "\n".join(document_chunks[i] for i in idxs[0] if i < len(document_chunks))


# Defining Funtion for prompts from the LLM model 

def ask_question(question):
    context = retrieve_context(question)
    prompt = f"""Answer the question based on the context below:

Context:
{context}

Question: {question}
Answer:"""

    response = client.chat.completions.create(
        model="llama3-70b-8192", #you can change the model according to your choice
        messages=[
            {"role": "system", "content": "Be short and precise."},
            {"role": "user", "content": prompt}
        ],
        temperature=0,
        max_tokens=500,
    )

    return response.choices[0].message.content.strip()


#Main Loop 

if __name__ == "__main__":
    file_path = input("Enter file path: ").strip()
    ext = os.path.splitext(file_path)[1].lower()

    loader_map = {
        ".pdf": load_pdf, ".docx": load_docx, ".pptx": load_pptx,
        ".xlsx": load_excel, ".csv": load_csv, ".json": load_json,
        ".txt": load_txt, ".png": load_image,
        ".jpg": load_image, ".jpeg": load_image
    }

    if ext not in loader_map:
        print(f"❌ Unsupported file type: {ext}")
        exit(1)

    raw = loader_map[ext](file_path)
    chunks = split_text(raw)
    add_documents(chunks)
    print("✅ Ready to chat!")


    while True:
        q = input("\nYou: ").strip()
        if q.lower() in ("exit", "quit"):
            print("👋 Goodbye!")
            break
        ans = ask_question(q)
        print("Bot:", ans)

# End Of Code